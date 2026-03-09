"""
Deck Generator — Build a clean, investor-ready PPTX pitch deck.

Uses python-pptx to create a professional 12-slide presentation with:
  - Consistent branding (gradient title slide, accent colours)
  - Proper slide layouts (title + body)
  - Section-by-section content injection

The generated file is saved to ``uploads/`` and returned as a download.
"""

from __future__ import annotations

import os
import logging
from datetime import datetime
from typing import Dict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------

CLR_PRIMARY   = RGBColor(0x4F, 0x46, 0xE5)   # Indigo 600 (matches Tailwind theme)
CLR_DARK      = RGBColor(0x1E, 0x1B, 0x4B)   # Indigo 950
CLR_LIGHT_BG  = RGBColor(0xF5, 0xF3, 0xFF)   # Very light indigo
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_TEXT      = RGBColor(0x1F, 0x2A, 0x37)    # Cool gray 800
CLR_SUBTEXT   = RGBColor(0x6B, 0x72, 0x80)    # Cool gray 500
CLR_ACCENT    = RGBColor(0x06, 0xB6, 0xD4)    # Cyan 500

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Canonical slide order
SLIDE_ORDER = [
    "Title",
    "Problem",
    "Solution",
    "Market",
    "Product",
    "Business Model",
    "Traction",
    "Competition",
    "Go-To-Market",
    "Team",
    "Financials",
    "Ask",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _add_solid_bg(slide, color: RGBColor) -> None:
    """Set a solid background colour on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = CLR_TEXT,
    alignment=PP_ALIGN.LEFT,
) -> None:
    """Helper to add a simple textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Split by newline → separate paragraphs for multi-line content
    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Detect bullet points (lines starting with - or •)
        if line.startswith(("- ", "• ", "* ")):
            line = line[2:].strip()
            p.level = 0
            p.space_before = Pt(4)

        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation, startup_name: str = "Your Startup") -> None:
    """Slide 1: Title slide with branded background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_solid_bg(slide, CLR_PRIMARY)

    _add_textbox(
        slide,
        left=Inches(1.5), top=Inches(2.2),
        width=Inches(10), height=Inches(1.5),
        text=startup_name,
        font_size=44, bold=True, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        left=Inches(1.5), top=Inches(3.8),
        width=Inches(10), height=Inches(1),
        text="Investor Pitch Deck",
        font_size=24, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        left=Inches(1.5), top=Inches(5.0),
        width=Inches(10), height=Inches(0.6),
        text=datetime.now().strftime("%B %Y"),
        font_size=16, color=CLR_WHITE,
        alignment=PP_ALIGN.CENTER,
    )


def _build_section_slide(
    prs: Presentation,
    title: str,
    content: str,
    slide_num: int,
) -> None:
    """Generic content slide: coloured header bar + body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_solid_bg(slide, CLR_WHITE)

    # Header bar (left accent strip)
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0), Inches(0),
        Inches(0.25), SLIDE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_PRIMARY
    shape.line.fill.background()

    # Slide number badge
    _add_textbox(
        slide,
        left=Inches(0.5), top=Inches(0.3),
        width=Inches(1), height=Inches(0.5),
        text=f"{slide_num:02d}",
        font_size=14, color=CLR_SUBTEXT,
    )

    # Section title
    _add_textbox(
        slide,
        left=Inches(0.8), top=Inches(0.6),
        width=Inches(11), height=Inches(1),
        text=title,
        font_size=32, bold=True, color=CLR_DARK,
    )

    # Divider line
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.5), Inches(11), Pt(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_PRIMARY
    line.line.fill.background()

    # Body content
    _add_textbox(
        slide,
        left=Inches(0.8), top=Inches(1.8),
        width=Inches(11), height=Inches(5),
        text=content,
        font_size=18, color=CLR_TEXT,
    )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def generate_pptx(
    improved_deck: Dict[str, str],
    output_dir: str = "uploads",
    startup_name: str = "Your Startup",
) -> str:
    """
    Build a clean 12-slide PPTX from improved section content.

    Args:
        improved_deck: Section name → polished text (from ``improve_deck``).
        output_dir: Directory to save the file.
        startup_name: Name shown on the title slide.

    Returns:
        Absolute path to the generated ``.pptx`` file.
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    _build_title_slide(prs, startup_name)

    # Slides 2-12: sections in canonical order
    slide_num = 2
    for section in SLIDE_ORDER[1:]:  # skip "Title"
        content = improved_deck.get(section, "(Content not available)")
        _build_section_slide(prs, section, content, slide_num)
        slide_num += 1

    # Save
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"Refined_Pitch_{timestamp}.pptx"
    file_path = os.path.join(output_dir, filename)
    prs.save(file_path)

    logger.info(f"[DECK-GEN] Generated refined deck: {file_path}")
    return os.path.abspath(file_path)
