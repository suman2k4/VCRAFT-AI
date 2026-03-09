"""
File Utilities — Slide extraction for PPTX and PDF pitch decks.

Handles:
  - PPTX → python-pptx: extract title + bullet text per slide
  - PDF  → pdfplumber: extract text per page
  - Text sanitisation (remove null bytes, excessive whitespace)
  - File-type & size validation helpers

Output format (shared):
  [
      {"slide_number": 1, "content": "..."},
      {"slide_number": 2, "content": "..."},
  ]
"""

from __future__ import annotations

import os
import re
import logging
from pathlib import Path
from typing import List, Dict

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_FILE_SIZE_MB = 10
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
ALLOWED_EXTENSIONS = {".pptx", ".pdf"}

# ---------------------------------------------------------------------------
# Validation helpers
# ---------------------------------------------------------------------------


def validate_file(filename: str, file_size: int) -> None:
    """
    Raise ``ValueError`` if the file is not an allowed type or exceeds size.

    Args:
        filename: Original filename (used for extension check).
        file_size: Size in bytes.

    Raises:
        ValueError: With a user-friendly message.
    """
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. Only .pptx and .pdf files are allowed."
        )
    if file_size > MAX_FILE_SIZE_BYTES:
        raise ValueError(
            f"File too large ({file_size / 1024 / 1024:.1f} MB). "
            f"Maximum allowed size is {MAX_FILE_SIZE_MB} MB."
        )


# ---------------------------------------------------------------------------
# Text sanitisation
# ---------------------------------------------------------------------------


def _sanitize(text: str) -> str:
    """Remove null bytes, collapse excessive whitespace, strip edges."""
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)            # horizontal whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)          # vertical whitespace
    return text.strip()


# ---------------------------------------------------------------------------
# PPTX Extraction
# ---------------------------------------------------------------------------


def extract_pptx(file_path: str) -> List[Dict]:
    """
    Extract slide content from a PowerPoint (.pptx) file.

    For each slide we look at:
      1. The slide title (if present)
      2. All text frames (bullet points, body text, etc.)

    Returns:
        List of dicts: ``[{"slide_number": 1, "content": "..."}, ...]``
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_data: List[Dict] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        # Try to grab the slide title shape first
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(slide.shapes.title.text.strip())

        # Walk every shape's text frame
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = paragraph.text.strip()
                if line and line not in parts:          # avoid dupe of title
                    parts.append(line)

        content = _sanitize("\n".join(parts))
        if content:
            slides_data.append({"slide_number": idx, "content": content})
        else:
            # Keep empty slides so numbering stays consistent
            slides_data.append({"slide_number": idx, "content": "(empty slide)"})

    logger.info(f"[FILE-UTILS] Extracted {len(slides_data)} slides from PPTX")
    return slides_data


# ---------------------------------------------------------------------------
# PDF Extraction
# ---------------------------------------------------------------------------


def extract_pdf(file_path: str) -> List[Dict]:
    """
    Extract text per page from a PDF file using pdfplumber.

    Falls back to PyPDF2 if pdfplumber produces empty results (e.g. scanned PDFs
    sometimes have invisible OCR text that PyPDF2 picks up).

    Returns:
        List of dicts: ``[{"slide_number": 1, "content": "..."}, ...]``
    """
    pages_data: List[Dict] = []

    # --- Primary: pdfplumber ---
    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                raw = page.extract_text() or ""
                content = _sanitize(raw)
                pages_data.append({
                    "slide_number": page_num,
                    "content": content if content else "(empty page)",
                })

        total_chars = sum(len(p["content"]) for p in pages_data)
        if total_chars > 50:
            logger.info(
                f"[FILE-UTILS] Extracted {len(pages_data)} pages from PDF "
                f"({total_chars} chars) via pdfplumber"
            )
            return pages_data
    except Exception as exc:
        logger.warning(f"[FILE-UTILS] pdfplumber failed, falling back to PyPDF2: {exc}")

    # --- Fallback: PyPDF2 ---
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    pages_data = []
    for page_num, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        content = _sanitize(raw)
        pages_data.append({
            "slide_number": page_num,
            "content": content if content else "(empty page)",
        })

    logger.info(f"[FILE-UTILS] Extracted {len(pages_data)} pages from PDF via PyPDF2")
    return pages_data


# ---------------------------------------------------------------------------
# Unified dispatcher
# ---------------------------------------------------------------------------


def extract_slides(file_path: str) -> List[Dict]:
    """
    Auto-detect file type and extract slides/pages.

    Args:
        file_path: Path to a .pptx or .pdf file on disk.

    Returns:
        Structured list of ``{"slide_number", "content"}``.

    Raises:
        ValueError: If the extension is not supported.
        FileNotFoundError: If the path doesn't exist.
    """
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".pdf":
        return extract_pdf(file_path)
    else:
        raise ValueError(f"Unsupported extension '{ext}'")
